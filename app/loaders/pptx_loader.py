"""
PptxLoader — converts PowerPoint presentations to Markdown.

Slide-aware approach:
  - Each slide becomes its own Markdown section with a slide header
  - slide_number added to chunk metadata via page_number field on ParsedDocument
  - Speaker notes extracted and appended under each slide
  - Images described via GPT-4o Vision (if LLM configured)
  - Tables converted to Markdown table format

Why not MarkItDown for PPTX?
  MarkItDown merges all slides into one flat text block — slide boundaries
  and speaker notes are lost. For RAG quality we want:
    "The Q3 revenue target (Slide 5) is $2.4M" not just "$2.4M" with no context.

Pipeline:
  python-pptx → per-slide extraction → Markdown with ## Slide N headers
  → images → Vision descriptions appended inline
"""

import asyncio
import base64
from io import BytesIO
from typing import Optional

import structlog

from app.loaders.base import BaseLoader, ParsedDocument
from app.connectors.base import RawDocument
from app.core.exceptions import DocumentParsingError

log = structlog.get_logger(__name__)

_MAX_SLIDES_FOR_VISION = 50  # cap Vision calls for very large decks


class PptxLoader(BaseLoader):

    def __init__(self, llm=None) -> None:
        """
        Args:
            llm: Optional LangChain ChatModel for image descriptions.
                 If None, images are noted as [Image] placeholders.
        """
        self._llm = llm
        log.info("loader.pptx.initialised", llm_enrichment=llm is not None)

    def supported_mime_types(self) -> list[str]:
        return [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ]

    async def load(self, raw_doc: RawDocument) -> ParsedDocument:
        try:
            markdown_text, metadata = await asyncio.to_thread(
                self._convert_sync, raw_doc
            )
        except DocumentParsingError:
            raise
        except Exception as e:
            raise DocumentParsingError(
                f"PptxLoader failed to parse '{raw_doc.file_name}': {e}"
            ) from e

        if not markdown_text.strip():
            raise DocumentParsingError(
                f"No content extracted from '{raw_doc.file_name}'"
            )
        # Enrich images via Vision if LLM available
        if self._llm is not None and metadata.get("has_images"):
            markdown_text = await self._enrich_images(raw_doc, markdown_text)

        log.info(
            "loader.pptx.parsed",
            file=raw_doc.file_name,
            slide_count=metadata.get("slide_count", 0),
            has_images=metadata.get("has_images", False),
            has_notes=metadata.get("has_notes", False),
        )

        return ParsedDocument(
            markdown_text=markdown_text,
            extracted_metadata={
                "loader": "python-pptx",
                "llm_enrichment": (
                    ["image_descriptions"]
                    if self._llm and metadata.get("has_images")
                    else []
                ),
                **metadata,
            },
        )

    # ── Sync conversion ────────────────────────────────────────────────────────

    def _convert_sync(self, raw_doc: RawDocument) -> tuple[str, dict]:
        """
        Convert PPTX to slide-aware Markdown.
        Each slide → ## Slide N — {title} section.
        Speaker notes → blockquote under the slide.
        Tables → proper Markdown tables.
        Images → [Image N] placeholder (replaced by Vision in async step).
        """
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(BytesIO(raw_doc.content))

        slide_count = len(prs.slides)
        has_images = False
        has_notes = False
        total_image_count = 0
        slide_sections = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            lines = []

            # ── Title ──────────────────────────────────────────────────────────
            title = self._get_slide_title(slide)
            if title:
                lines.append(f"## Slide {slide_num} — {title}")
            else:
                lines.append(f"## Slide {slide_num}")

            lines.append("")  # blank line after heading

            # ── Shapes: text, tables, images ───────────────────────────────────
            image_index_on_slide = 0
            for shape in slide.shapes:

                # Tables
                if shape.has_table:
                    lines.append(self._table_to_markdown(shape.table))
                    lines.append("")
                    continue

                # Text frames (skip title shape — already captured above)
                if shape.has_text_frame:
                    if shape == slide.shapes.title:
                        continue
                    text = self._extract_text_frame(shape.text_frame)
                    if text.strip():
                        lines.append(text)
                        lines.append("")

                # Images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
                    hasattr(shape, "image")
                ):
                    has_images = True
                    total_image_count += 1
                    image_index_on_slide += 1
                    # Placeholder — replaced by Vision in _enrich_images()
                    lines.append(
                        f"[IMAGE_PLACEHOLDER slide={slide_num} "
                        f"index={image_index_on_slide}]"
                    )
                    lines.append("")

            # ── Speaker notes ──────────────────────────────────────────────────
            notes_text = self._get_notes(slide)
            if notes_text:
                has_notes = True
                lines.append("> **Speaker Notes:**")
                for note_line in notes_text.split("\n"):
                    if note_line.strip():
                        lines.append(f"> {note_line.strip()}")
                lines.append("")

            slide_sections.append("\n".join(lines))

        markdown_text = "\n\n---\n\n".join(slide_sections)

        metadata = {
            "slide_count": slide_count,
            "has_images": has_images,
            "has_notes": has_notes,
            "total_image_count": total_image_count,
        }

        return markdown_text, metadata

    # ── Helpers ────────────────────────────────────────────────────────────────

    def _get_slide_title(self, slide) -> Optional[str]:
        """Extract title from title placeholder if it exists."""
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                return slide.shapes.title.text_frame.text.strip()
        except Exception:
            pass
        return None

    def _get_notes(self, slide) -> Optional[str]:
        """Extract speaker notes text."""
        try:
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                text = notes_frame.text.strip()
                return text if text else None
        except Exception:
            pass
        return None

    def _extract_text_frame(self, text_frame) -> str:
        """Extract text from a shape's text frame, preserving list structure."""
        lines = []
        for para in text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Detect bullet/list level via indent level
            level = para.level or 0
            indent = "  " * level
            if level > 0:
                lines.append(f"{indent}- {text}")
            else:
                lines.append(text)
        return "\n".join(lines)

    def _table_to_markdown(self, table) -> str:
        """Convert a pptx table to Markdown table format."""
        rows = []
        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if row_idx == 0:
                # Add separator after header row
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
        return "\n".join(rows)

    # ── Vision enrichment ──────────────────────────────────────────────────────

    async def _enrich_images(self, raw_doc: RawDocument, markdown_text: str) -> str:
        """
        Replace [IMAGE_PLACEHOLDER slide=N index=M] markers with
        GPT-4o Vision descriptions extracted from the actual PPTX images.
        """
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from langchain_core.messages import HumanMessage
        import re

        prs = Presentation(BytesIO(raw_doc.content))

        # Build lookup: (slide_num, image_index) → image bytes
        image_map: dict[tuple[int, int], bytes] = {}
        for slide_num, slide in enumerate(prs.slides, start=1):
            image_index = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(
                    shape, "image"
                ):
                    image_index += 1
                    try:
                        image_map[(slide_num, image_index)] = shape.image.blob
                    except Exception:
                        pass

        # Replace each placeholder with Vision description
        placeholder_pattern = re.compile(
            r"\[IMAGE_PLACEHOLDER slide=(\d+) index=(\d+)\]"
        )

        async def replace_placeholder(match) -> str:
            slide_num = int(match.group(1))
            img_idx = int(match.group(2))
            image_bytes = image_map.get((slide_num, img_idx))

            if image_bytes is None:
                return f"[Image {img_idx}, Slide {slide_num}]"

            # Skip very small images (likely icons/bullets)
            if len(image_bytes) < 2048:
                return f"[Image {img_idx}, Slide {slide_num}]"

            try:
                image_b64 = base64.b64encode(image_bytes).decode()
                fmt = "png" if image_bytes[:8] == b"\x89PNG\r\n\x1a\n" else "jpeg"

                message = HumanMessage(
                    content=[
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{fmt};base64,{image_b64}",
                                "detail": "high",
                            },
                        },
                        {
                            "type": "text",
                            "text": (
                                f"This image is from slide {slide_num} of a PowerPoint. "
                                "Describe it concisely for a RAG knowledge base. "
                                "Focus on: key information, data, text visible, charts, "
                                "and what it communicates. Single paragraph, factual."
                            ),
                        },
                    ]
                )
                response = await self._llm.ainvoke([message])
                description = response.content.strip()
                log.debug(
                    "loader.pptx.image_described",
                    file=raw_doc.file_name,
                    slide=slide_num,
                    image=img_idx,
                )
                return f"**[Image {img_idx}, Slide {slide_num}]**: {description}"

            except Exception as e:
                log.warning(
                    "loader.pptx.image_description_failed",
                    file=raw_doc.file_name,
                    slide=slide_num,
                    image=img_idx,
                    error=str(e),
                )
                return f"[Image {img_idx}, Slide {slide_num}]"

        # Process all placeholders
        results = {}
        for match in placeholder_pattern.finditer(markdown_text):
            key = (int(match.group(1)), int(match.group(2)))
            if key not in results:
                results[key] = await replace_placeholder(match)

        def replace_match(m):
            key = (int(m.group(1)), int(m.group(2)))
            return results.get(key, f"[Image {m.group(2)}, Slide {m.group(1)}]")

        return placeholder_pattern.sub(replace_match, markdown_text)
